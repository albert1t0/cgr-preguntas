from flask import Flask, render_template, request, redirect, url_for, send_file
import os
from werkzeug.utils import secure_filename
import PyPDF2
import pptx
import google.generativeai as genai
import json
from dotenv import load_dotenv
import base64
import mimetypes

load_dotenv()

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = "uploads"

# Configure the Gemini API key
try:
    api_key = os.environ.get("GOOGLE_API_KEY")
    if not api_key:
        raise ValueError(
            "La variable de entorno GOOGLE_API_KEY no está configurada o el archivo .env no se encontró."
        )
    genai.configure(api_key=api_key)
except (KeyError, ValueError) as e:
    # This will stop the app from starting if the key is not found
    raise SystemExit(e)


def process_background_image(image_file):
    """Process background image and convert to base64"""
    if not image_file or image_file.filename == "":
        return None

    try:
        # Read the image file
        image_data = image_file.read()
        image_file.seek(0)  # Reset file pointer

        # Get MIME type
        mime_type, _ = mimetypes.guess_type(image_file.filename)
        if not mime_type or not mime_type.startswith("image/"):
            return None

        # Convert to base64
        base64_image = base64.b64encode(image_data).decode("utf-8")
        return f"data:{mime_type};base64,{base64_image}"

    except Exception as e:
        print(f"Error procesando imagen de fondo: {e}")
        return None


def extract_text(filepath):
    ext = filepath.rsplit(".", 1)[1].lower()
    text = ""
    if ext == "pdf":
        try:
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text()
        except PyPDF2.errors.PdfReadError:
            return "Error: No se pudo leer el archivo PDF. Puede que esté corrupto o encriptado."
    elif ext == "pptx":
        try:
            prs = pptx.Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            return f"Error: No se pudo leer el archivo PPTX. {e}"
    return text


def generate_quiz(text, num_questions=2):
    model = genai.GenerativeModel("gemini-1.5-flash")

    # Generate example structure based on number of questions
    example_questions = []
    for i in range(num_questions):
        example_questions.append(
            {
                "pregunta": f"¿Tu pregunta {i + 1} aquí?",
                "opciones": ["Opción A", "Opción B", "Opción C", "Opción D"],
                "respuesta_correcta_index": i % 4,
                "justificacion": f"La justificación de la pregunta {i + 1}.",
            }
        )

    example_structure = {"preguntas": example_questions}

    prompt = (
        f"""
    Basándote en el siguiente texto, genera un cuestionario de {num_questions} preguntas de opción múltiple.
    El formato de salida debe ser un objeto JSON válido y nada más. No incluyas '```json' o '```' en la salida.
    La estructura debe ser la siguiente:
    {json.dumps(example_structure, indent=2, ensure_ascii=False)}

    Texto de contexto: Las respuesta de justificación deben iniciar por "Según al información dsarrollada en el tema ..." o "Basando en el tema desarrollado ...", simepre refiriendose al tema.
    """
        + text
    )

    try:
        response = model.generate_content(prompt)
        # Clean the response to get only the JSON part
        json_response_cleaned = (
            response.text.strip().replace("```json", "").replace("```", "")
        )
        quiz = json.loads(json_response_cleaned)
        return quiz
    except Exception as e:
        # It's better to return the error and handle it in the route
        return {
            "error": f"Error al generar o parsear el cuestionario desde Gemini: {e}",
            "raw_response": response.text if "response" in locals() else "No response",
        }


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return redirect(request.url)
    file = request.files["file"]
    if file.filename == "":
        return redirect(request.url)
    if file:
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(filepath)

        # Get number of questions from form
        num_questions = int(request.form.get("num_questions", 2))

        # Process background image if provided
        background_image = None
        if "background_image" in request.files:
            background_image = process_background_image(
                request.files["background_image"]
            )

        extracted_text = extract_text(filepath)
        if extracted_text.startswith("Error:"):
            return (
                f"<h1>Error en la extracción de texto</h1><p>{extracted_text}</p>",
                400,
            )

        quiz_data = generate_quiz(extracted_text, num_questions)

        if "error" in quiz_data:
            return (
                f"<h1>Error al generar el Quiz</h1><pre>{json.dumps(quiz_data, indent=2)}</pre>",
                500,
            )

        # Render the quiz template with the data
        rendered_html = render_template(
            "quiz.html", quiz_data=quiz_data, background_image=background_image
        )

        # Save the rendered HTML to a file
        output_filename = f"test_{os.path.splitext(filename)[0]}.html"
        output_filepath = os.path.join(app.config["UPLOAD_FOLDER"], output_filename)

        with open(output_filepath, "w", encoding="utf-8") as f:
            f.write(rendered_html)

        # Send the generated file to the user for download
        return send_file(output_filepath, as_attachment=True)


if __name__ == "__main__":
    os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
    app.run(debug=True, port=5001)
